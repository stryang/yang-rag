"""Document loader for multiple file formats - using direct pypdf/text APIs."""

from pathlib import Path
from typing import List, Optional

from langchain_core.documents import Document


class MultiFormatLoader:
    """Unified loader that detects file format and uses appropriate loader."""

    LOADER_MAP = {
        ".pdf": "load_pdf",
        ".txt": "load_text",
        ".md": "load_text",
        ".markdown": "load_text",
        ".html": "load_text",
        ".htm": "load_text",
        ".docx": "load_docx",
        ".doc": "load_docx",
        ".pptx": "load_pptx",
        ".ppt": "load_pptx",
        ".xlsx": "load_text",
        ".xls": "load_text",
    }

    def __init__(self):
        """Initialize loader."""
        pass

    def load(self, file_path: str) -> List[Document]:
        """Load document with automatic format detection.

        Args:
            file_path: Path to the document file

        Returns:
            List of Document objects
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError("File not found: {}".format(file_path))

        suffix = file_path.suffix.lower()
        loader_method = self.LOADER_MAP.get(suffix)

        if loader_method is None:
            raise ValueError("Unsupported file type: {}".format(suffix))

        docs = getattr(self, loader_method)(str(file_path))

        for doc in docs:
            doc.metadata["source"] = str(file_path)
            doc.metadata["file_name"] = file_path.name
            doc.metadata["file_type"] = suffix

        return docs

    def load_pdf(self, file_path: str) -> List[Document]:
        """Load PDF file using pypdf directly."""
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        docs = []

        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"page": page_num + 1}
                ))

        return docs

    def load_text(self, file_path: str) -> List[Document]:
        """Load text/markdown/HTML file."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        return [Document(page_content=content, metadata={})]

    def load_docx(self, file_path: str) -> List[Document]:
        """Load Word document using python-docx."""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

            if paragraphs:
                content = "\n\n".join(paragraphs)
                return [Document(page_content=content, metadata={})]
            return []
        except ImportError:
            return self.load_text(file_path)

    def load_pptx(self, file_path: str) -> List[Document]:
        """Load PowerPoint file using python-pptx."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides_content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    slides_content.append("Slide {}:\n{}".format(
                        slide_num, "\n".join(slide_text)
                    ))

            if slides_content:
                content = "\n\n".join(slides_content)
                return [Document(page_content=content, metadata={})]
            return []
        except ImportError:
            return []

    def load_bytes(
        self, content: bytes, file_name: str, file_type: Optional[str] = None
    ) -> List[Document]:
        """Load document from bytes content."""
        import tempfile

        file_type = file_type or Path(file_name).suffix.lower()
        suffix = file_type if file_type.startswith(".") else ".{}".format(file_type)

        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            return self.load(tmp_path)
        finally:
            Path(tmp_path).unlink(missing_ok=True)


_loader_instance = None


def get_loader() -> MultiFormatLoader:
    """Get singleton loader instance."""
    global _loader_instance
    if _loader_instance is None:
        _loader_instance = MultiFormatLoader()
    return _loader_instance
